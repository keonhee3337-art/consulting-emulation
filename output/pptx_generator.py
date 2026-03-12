"""
pptx_generator.py — Step 3.1 of Automated M&A Due Diligence & Strategy War Room

Converts the FinAgent Report Agent's markdown output into a consulting-style
5-slide PowerPoint deck using python-pptx.

Slide structure:
  1. Title slide (navy bg)
  2. Key Findings (white bg)
  3. Financial Summary (white bg, table if markdown table present)
  4. Market & Strategic Context (grey bg)
  5. Analyst Note + optional Valuation result (white bg)

Usage:
  from output.pptx_generator import generate_deck
  path = generate_deck(report_text, "Samsung Electronics", "output/deck.pptx")

Standalone test:
  python output/pptx_generator.py
"""

import re
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)   # title slide background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # body background / white text
GREY   = RGBColor(0xF5, 0xF5, 0xF5)   # accent slide background
GOLD   = RGBColor(0xC9, 0xA8, 0x4C)   # accent / rule line
DARK   = RGBColor(0x1A, 0x1A, 0x1A)   # body text

# ---------------------------------------------------------------------------
# Slide dimensions — standard 16:9 widescreen
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Font sizes
# ---------------------------------------------------------------------------
TITLE_PT = 32
BODY_PT  = 18
SMALL_PT = 12


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb_fill(shape, color: RGBColor):
    """Fill a shape's background with a solid colour."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _slide_bg(slide, color: RGBColor):
    """Set the entire slide background to a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text: str, left, top, width, height,
                 font_size=BODY_PT, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a plain textbox and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def _add_title_bar(slide, title_text: str, bg_color=NAVY, text_color=WHITE):
    """Add a full-width title bar at the top of the slide."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        SLIDE_W, Inches(1.1)
    )
    _rgb_fill(bar, bg_color)
    bar.line.fill.background()  # no border

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(TITLE_PT)
    run.font.bold = True
    run.font.color.rgb = text_color

    # Left padding — shift text right via a narrow left margin shape
    bar.text_frame.margin_left = Inches(0.4)
    bar.text_frame.margin_top = Inches(0.15)


def _bullets_from_text(text: str, max_bullets: int = 6) -> list[str]:
    """
    Extract bullet points from a markdown section body.
    Handles lines starting with '- ', '* ', or plain text lines.
    Strips empty lines. Returns at most max_bullets items.
    """
    lines = text.strip().splitlines()
    bullets = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        # Strip leading markdown bullet markers
        line = re.sub(r'^[-*]\s+', '', line)
        if line:
            bullets.append(line)
        if len(bullets) >= max_bullets:
            break
    return bullets


def _add_bullet_list(slide, bullets: list[str],
                     top=Inches(1.3), max_h=Inches(5.8)):
    """
    Add a vertical list of bullet items to the slide body area.
    Each bullet gets its own text run inside a single text frame.
    """
    if not bullets:
        _add_textbox(slide, "No data available.",
                     Inches(0.5), top, Inches(12.3), Inches(1))
        return

    txBox = slide.shapes.add_textbox(
        Inches(0.5), top, Inches(12.3), max_h
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"  \u2022  {bullet}"
        run.font.size = Pt(BODY_PT)
        run.font.color.rgb = DARK


# ---------------------------------------------------------------------------
# Markdown section parser
# ---------------------------------------------------------------------------

def parse_markdown_sections(report_text: str) -> dict:
    """
    Extract the four sections produced by the FinAgent Report Agent:
      - Key Findings
      - Financial Data Summary
      - Market & Strategic Context
      - Analyst Note

    Returns a dict { section_name: body_text }.
    If a section is missing, its value is an empty string.
    """
    section_names = [
        "Key Findings",
        "Financial Data Summary",
        "Market & Strategic Context",
        "Analyst Note",
    ]

    sections: dict[str, str] = {name: "" for name in section_names}

    # Split on ## headings; keep heading in token
    parts = re.split(r'(##\s+.+)', report_text)

    current_section = None
    for part in parts:
        header_match = re.match(r'##\s+(.+)', part.strip())
        if header_match:
            heading = header_match.group(1).strip()
            # Match to known section names (partial match for flexibility)
            current_section = None
            for name in section_names:
                if name.lower() in heading.lower():
                    current_section = name
                    break
        elif current_section is not None:
            sections[current_section] += part

    # Strip leading/trailing whitespace from each section body
    return {k: v.strip() for k, v in sections.items()}


# ---------------------------------------------------------------------------
# Markdown table parser
# ---------------------------------------------------------------------------

def _parse_markdown_table(text: str) -> tuple[list[str], list[list[str]]] | None:
    """
    Detect and parse a markdown table in the text.
    Returns (headers, rows) or None if no table found.

    Expects format:
      | Col A | Col B |
      |-------|-------|
      | val1  | val2  |
    """
    lines = text.strip().splitlines()
    table_lines = [l for l in lines if l.strip().startswith('|')]

    if len(table_lines) < 3:
        return None

    def split_row(line):
        return [cell.strip() for cell in line.strip().strip('|').split('|')]

    headers = split_row(table_lines[0])
    # table_lines[1] is the separator row --- skip it
    rows = [split_row(l) for l in table_lines[2:] if '---' not in l]

    return headers, rows


def _add_pptx_table(slide, headers: list[str], rows: list[list[str]],
                    top=Inches(1.3)):
    """
    Render a simple python-pptx table from headers + rows.
    Max 8 rows to keep slide readable.
    """
    rows = rows[:8]
    n_rows = len(rows) + 1   # +1 for header
    n_cols = len(headers)

    col_w = Inches(12.0 / n_cols)
    row_h = Inches(0.45)

    table = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), top,
        Inches(12.0), Inches(row_h.inches * n_rows)
    ).table

    # Style header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.size = Pt(SMALL_PT + 2)
        run.font.color.rgb = WHITE

    # Style data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            if col_idx >= n_cols:
                continue
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
            # Alternate row shading
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEA, 0xEF, 0xF5)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(SMALL_PT)
            run.font.color.rgb = DARK


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _slide_1_title(prs: Presentation, company_name: str):
    """Slide 1: Title — dark navy with white text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _slide_bg(slide, NAVY)

    # Gold horizontal rule
    rule = slide.shapes.add_shape(
        1,
        Inches(0.6), Inches(3.5),
        Inches(5.5), Pt(3)
    )
    _rgb_fill(rule, GOLD)
    rule.line.fill.background()

    # Main title
    _add_textbox(
        slide,
        f"M&A Due Diligence:\n{company_name}",
        Inches(0.6), Inches(2.0),
        Inches(11.0), Inches(1.8),
        font_size=40, bold=True, color=WHITE,
        align=PP_ALIGN.LEFT
    )

    # Subtitle
    _add_textbox(
        slide,
        f"Strategy War Room  |  {date.today().strftime('%B %d, %Y')}",
        Inches(0.6), Inches(3.8),
        Inches(10.0), Inches(0.6),
        font_size=BODY_PT, bold=False, color=GOLD,
        align=PP_ALIGN.LEFT
    )

    # Confidentiality notice — bottom left
    _add_textbox(
        slide,
        "Confidential — Internal Use Only",
        Inches(0.4), Inches(6.9),
        Inches(6.0), Inches(0.4),
        font_size=SMALL_PT, bold=False,
        color=RGBColor(0xAA, 0xAA, 0xAA),
        align=PP_ALIGN.LEFT
    )


def _slide_2_key_findings(prs: Presentation, section_text: str):
    """Slide 2: Key Findings — white bg, bullet list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, WHITE)
    _add_title_bar(slide, "Key Findings", bg_color=NAVY, text_color=WHITE)

    bullets = _bullets_from_text(section_text, max_bullets=6)
    _add_bullet_list(slide, bullets, top=Inches(1.3))

    # Gold accent rule below title bar
    rule = slide.shapes.add_shape(
        1, Inches(0), Inches(1.1), SLIDE_W, Pt(3)
    )
    _rgb_fill(rule, GOLD)
    rule.line.fill.background()


def _slide_3_financial(prs: Presentation, section_text: str):
    """Slide 3: Financial Summary — renders markdown table if present."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, WHITE)
    _add_title_bar(slide, "Financial Summary", bg_color=NAVY, text_color=WHITE)

    # Gold rule
    rule = slide.shapes.add_shape(
        1, Inches(0), Inches(1.1), SLIDE_W, Pt(3)
    )
    _rgb_fill(rule, GOLD)
    rule.line.fill.background()

    if not section_text:
        _add_textbox(slide, "No data available.",
                     Inches(0.5), Inches(1.4), Inches(12.3), Inches(1))
        return

    table_result = _parse_markdown_table(section_text)
    if table_result:
        headers, rows = table_result
        _add_pptx_table(slide, headers, rows, top=Inches(1.4))
    else:
        # Fallback: render as bullets
        bullets = _bullets_from_text(section_text, max_bullets=8)
        _add_bullet_list(slide, bullets, top=Inches(1.3))


def _slide_4_strategic(prs: Presentation, section_text: str):
    """Slide 4: Market & Strategic Context — light grey bg."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, GREY)
    _add_title_bar(slide, "Market & Strategic Context",
                   bg_color=NAVY, text_color=WHITE)

    # Gold rule
    rule = slide.shapes.add_shape(
        1, Inches(0), Inches(1.1), SLIDE_W, Pt(3)
    )
    _rgb_fill(rule, GOLD)
    rule.line.fill.background()

    bullets = _bullets_from_text(section_text, max_bullets=6)
    _add_bullet_list(slide, bullets, top=Inches(1.3))


def _slide_5_analyst(prs: Presentation, section_text: str,
                     valuation_result: str | None):
    """Slide 5: Analyst Note + optional Valuation summary — white bg."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(slide, WHITE)
    _add_title_bar(slide, "Analyst Note", bg_color=NAVY, text_color=WHITE)

    # Gold rule
    rule = slide.shapes.add_shape(
        1, Inches(0), Inches(1.1), SLIDE_W, Pt(3)
    )
    _rgb_fill(rule, GOLD)
    rule.line.fill.background()

    analyst_text = section_text.strip() if section_text else "No analyst note available."
    # Remove leading bullet markers for the note (it's usually a single sentence)
    analyst_text = re.sub(r'^[-*]\s+', '', analyst_text)

    _add_textbox(
        slide, analyst_text,
        Inches(0.5), Inches(1.4),
        Inches(12.3), Inches(2.5),
        font_size=BODY_PT, color=DARK
    )

    if valuation_result:
        # Section label
        _add_textbox(
            slide, "Valuation Summary",
            Inches(0.5), Inches(4.0),
            Inches(12.3), Inches(0.4),
            font_size=SMALL_PT + 2, bold=True, color=NAVY
        )
        # Thin gold divider
        div = slide.shapes.add_shape(
            1, Inches(0.5), Inches(4.45), Inches(12.0), Pt(1)
        )
        _rgb_fill(div, GOLD)
        div.line.fill.background()

        # Valuation content at smaller font
        val_text = valuation_result.strip()
        _add_textbox(
            slide, val_text,
            Inches(0.5), Inches(4.6),
            Inches(12.3), Inches(2.5),
            font_size=SMALL_PT, color=DARK
        )


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def generate_deck(
    report: str,
    company_name: str,
    output_path: str,
    valuation_result: str = None
) -> str:
    """
    Generate a consulting-style .pptx from the Report Agent's markdown output.

    Args:
        report:           Full markdown report string from report_agent.py
        company_name:     Target company name for title slide
        output_path:      Absolute or relative path for the .pptx file
        valuation_result: Optional plain-text DCF/comps summary from Valuation Agent

    Returns:
        output_path (the saved file location)
    """
    sections = parse_markdown_sections(report)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_1_title(prs, company_name)
    _slide_2_key_findings(prs, sections["Key Findings"])
    _slide_3_financial(prs, sections["Financial Data Summary"])
    _slide_4_strategic(prs, sections["Market & Strategic Context"])
    _slide_5_analyst(prs, sections["Analyst Note"], valuation_result)

    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# Standalone test — mock Samsung report matching Report Agent format
# ---------------------------------------------------------------------------

MOCK_REPORT = """
## Key Findings
- Samsung Electronics reported KRW 302.2T in revenue for FY2024, up 12% YoY driven by HBM3E ramp
- Operating profit recovered to KRW 32.7T (OPM 10.8%) after the memory downturn trough in 2023
- HBM3E yield issues in H1 2024 caused a 6-week supply delay to key US hyperscaler customers
- DRAM market share held at 43% globally; NAND share edged down to 31% amid SK Hynix pricing pressure
- Foundry (TSMC competitor) segment remains sub-scale: 3nm yield at ~55% vs TSMC's 72%
- Net cash position of KRW 98.4T provides capacity for a potential target acquisition or buyback

## Financial Data Summary
| Metric | FY2022 | FY2023 | FY2024 |
|--------|--------|--------|--------|
| Revenue (KRW T) | 302.2 | 258.9 | 302.2 |
| Operating Profit (KRW T) | 43.4 | 6.6 | 32.7 |
| OPM (%) | 14.4% | 2.5% | 10.8% |
| Net Income (KRW T) | 55.6 | 15.5 | 28.9 |
| CAPEX (KRW T) | 53.1 | 47.7 | 51.3 |
| Net Cash (KRW T) | 112.8 | 95.1 | 98.4 |

## Market & Strategic Context
- HBM (High Bandwidth Memory) is the fastest-growing segment in semiconductor; SK Hynix leads with HBM3E; Samsung is 6-9 months behind
- AI accelerator demand from NVIDIA, Google TPU, and AWS Trainium is the structural tailwind for premium DRAM
- TSMC's CoWoS advanced packaging creates bottleneck; Samsung Foundry's 2.5D packaging capability is a strategic differentiator if yield improves
- South Korea government chipact (K-Chips Act) provides 15-25% tax credit on facility investment — material for CAPEX-heavy fab builds
- Potential M&A targets: Netherlands-based ASML reseller, advanced packaging IP, or US-based HPC software layer to close the full-stack gap
- Geopolitical risk: US-China export controls limit Samsung's ability to expand leading-edge capacity in Xi'an facility

## Analyst Note
- Data sourced from DART disclosures and Samsung investor relations; HBM yield figures are consensus analyst estimates and should be validated against Samsung's next earnings call before use in client deliverables.
"""

MOCK_VALUATION = """DCF (base case): Intrinsic value KRW 81,400/share (current: KRW 74,200; upside 9.7%)
  WACC: 9.2%  |  Terminal growth: 3.0%  |  FCF CAGR (5yr): 11.4%
EV/EBITDA comps: Peer median 8.4x  |  Samsung implied 7.1x  →  ~18% discount to peers
Bear / Base / Bull: KRW 58,000 / KRW 81,400 / KRW 107,000"""

if __name__ == "__main__":
    import os

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_file = os.path.join(output_dir, "test_deck.pptx")

    path = generate_deck(
        report=MOCK_REPORT,
        company_name="Samsung Electronics",
        output_path=output_file,
        valuation_result=MOCK_VALUATION
    )
    print(f"Deck generated: {path}")
